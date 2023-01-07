import os
import io
import json
import boto3
import collections
import collections.abc
from pptx import Presentation
import nltk
from nltk.corpus import stopwords

# convertPPTXToText 
def lambda_handler(event, context):
    print('received event:')
    print(event)

    # Static Variables
    s3_client = boto3.client('s3')
    S3_BUCKET = os.environ['S3_BUCKET']

    db_client = boto3.client('dynamodb')
    DYNAMODB_TABLE = os.environ['DYNAMODB_TABLE']

    sqs_client = boto3.client('sqs')
    SQS_QUEUE = os.environ['SQS_QUEUE']

    id = event['Records'][0]['body']
    
    # Get s3key from DynamoDB
    response = db_client.get_item(
      TableName=DYNAMODB_TABLE,
      Key={
        'id':{'S': id}
      }
    )

    s3key = response['Item']['s3key']['S']

    # Extract Text
    response = s3_client.get_object(Bucket=S3_BUCKET, Key=s3key)
    prs_bytes = response['Body'].read()
    prs_file = io.BytesIO(prs_bytes)

    prs = Presentation(prs_file)

    text_runs = []

    for slide in prs.slides:
      for shape in slide.shapes:
          if not shape.has_text_frame:
              continue
          for paragraph in shape.text_frame.paragraphs:
              for run in paragraph.runs:
                  text_runs.append(run.text)

    raw_text = ' '.join(text_runs)
    
    # Remove stop words
    nltk.data.path.append("/tmp/nltk_data")
    nltk.download('stopwords', download_dir="/tmp/nltk_data")
    stopWords = stopwords.words('english')
    raw_text = ' '.join([word for word in raw_text.split() if word not in stopWords])

    # Update DynamoDB Entry
    response = db_client.update_item(
      TableName=DYNAMODB_TABLE,
      Key={
        'id':{'S': id}
      },
      UpdateExpression='SET raw_text = :raw_text',
      ExpressionAttributeValues={':raw_text': {'S': raw_text}}
    )
    
    # Create SQS Message
    response = sqs_client.send_message(
      QueueUrl=SQS_QUEUE,
      MessageBody=id
    )
